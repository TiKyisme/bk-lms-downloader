from __future__ import annotations

import importlib
import importlib.util
import io
import sys
import tempfile
import traceback
import zipfile
from contextlib import redirect_stdout
from dataclasses import dataclass
from pathlib import Path
from types import ModuleType, SimpleNamespace
from typing import Callable, Iterable

from .app_logging import get_logger
from .ai_study_pack import validate_ai_study_pack
from .models import Course


LOG = get_logger(__name__)
AI_TOOL_RELATIVE_PATH = Path("tools") / "prepare_ai_course.py"
REQUIRED_AI_MODULES = {
    "beautifulsoup4": "bs4",
    "markdownify": "markdownify",
    "pypdf": "pypdf",
    "python-pptx": "pptx",
}


class OptionalAIDependenciesError(RuntimeError):
    """Raised only for an incomplete developer/source installation."""


class AIPreparationError(RuntimeError):
    pass


def missing_ai_dependencies(
    importer: Callable[[str], object] = importlib.import_module,
) -> list[str]:
    """Return modules that cannot actually import in this runtime.

    ``find_spec`` is deliberately not used here: a frozen PyInstaller process
    can import a bundled module even when metadata/spec probing is inconsistent.
    """
    missing: list[str] = []
    for package, module in REQUIRED_AI_MODULES.items():
        try:
            importer(module)
        except ImportError:
            missing.append(package)
    return missing


def default_ai_output(course_root: Path) -> Path:
    return course_root / "AI_Knowledge"


def default_ai_tool_path() -> Path:
    """Find the bundled tool in PyInstaller or the checkout in source mode."""
    if getattr(sys, "frozen", False):
        root = Path(getattr(sys, "_MEIPASS", Path(sys.executable).parent))
    else:
        root = Path(__file__).resolve().parents[2]
    return root / AI_TOOL_RELATIVE_PATH


def ai_runtime_diagnostics() -> str:
    """Collect import facts for support without using metadata to gate runtime."""
    lines = [
        f"Frozen: {'yes' if getattr(sys, 'frozen', False) else 'no'}",
        f"sys._MEIPASS: {getattr(sys, '_MEIPASS', None)!r}",
        f"sys.executable: {sys.executable}",
        f"AI tool: {default_ai_tool_path()} ({'OK' if default_ai_tool_path().is_file() else 'MISSING'})",
    ]
    for package, module_name in REQUIRED_AI_MODULES.items():
        try:
            module = importlib.import_module(module_name)
        except Exception as exc:
            lines.append(f"{package} ({module_name}): IMPORT FAILED: {type(exc).__name__}: {exc}")
            continue
        lines.extend(
            (
                f"{package} ({module_name}): IMPORT OK",
                f"  __file__: {getattr(module, '__file__', None)!r}",
                f"  __spec__: {getattr(module, '__spec__', None)!r}",
            )
        )
    return "\n".join(lines)


def _ai_runtime_self_test() -> Path:
    """Exercise the same batch path used by the end-user GUI."""
    missing = missing_ai_dependencies()
    if missing:
        raise RuntimeError("Missing packaged AI modules: " + ", ".join(missing))
    if not default_ai_tool_path().is_file():
        raise RuntimeError("Bundled AI preparation tool is missing")

    with tempfile.TemporaryDirectory(prefix="bklms_ai_smoke_") as temp_dir:
        course_root = Path(temp_dir) / "Course"
        course_root.mkdir()
        (course_root / "notes.txt").write_text("AI runtime smoke", encoding="utf-8")

        html_dir = course_root / "Web Page"
        html_dir.mkdir()
        (html_dir / "content.html").write_text(
            "<h1>Packaged HTML smoke</h1><p>Local extraction.</p>",
            encoding="utf-8",
        )

        from pypdf import PdfWriter

        pdf_writer = PdfWriter()
        pdf_writer.add_blank_page(width=72, height=72)
        with (course_root / "tiny.pdf").open("wb") as pdf_handle:
            pdf_writer.write(pdf_handle)

        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Packaged PPTX smoke"
        presentation.save(course_root / "tiny.pptx")

        course = Course(
            id="ai-self-test",
            url="https://lms.hcmut.edu.vn/course/view.php?id=1",
            output=str(course_root),
            name="AI runtime self-test",
        )
        batch = AIBatchPreparer().prepare_courses(
            [course],
            lambda item: item.output_path,
        )
        if len(batch.succeeded) != 1 or batch.failed:
            detail = batch.failed[0].error if batch.failed else "unknown batch failure"
            raise RuntimeError("AI batch self-test failed: " + (detail or "unknown error"))
        output = batch.succeeded[0].output
        if output is None:
            raise RuntimeError("AI batch self-test returned no output")
        required_outputs = (
            output / "START_HERE.md",
            output / "COURSE_MAP.md",
            output / "COVERAGE_REPORT.md",
            output / "TUTOR_PROTOCOL.md",
            output / "CHATGPT_START_PROMPT.txt",
            output / "AI_TUTOR_CONTEXT.md",
            output / "course_index.md",
            output / "processing_report.md",
            output / "meta" / "corpus.jsonl",
        )
        if not all(path.is_file() for path in required_outputs):
            raise RuntimeError("AI runtime self-test did not create required outputs")
        validation = validate_ai_study_pack(output)
        if validation.errors:
            raise RuntimeError("AI Study Pack validation failed: " + "; ".join(validation.errors))
        packs = list(output.parent.glob("* - AI Study Pack.zip"))
        if len(packs) != 1:
            raise RuntimeError("AI runtime self-test did not create one ChatGPT-ready ZIP")
        with tempfile.TemporaryDirectory(prefix="bklms_ai_pack_roundtrip_") as unpacked:
            with zipfile.ZipFile(packs[0]) as archive:
                archive.extractall(unpacked)
            unpacked_validation = validate_ai_study_pack(Path(unpacked))
        if unpacked_validation.errors:
            raise RuntimeError(
                "AI Study Pack ZIP round-trip failed: " + "; ".join(unpacked_validation.errors)
            )
        return output


def run_ai_runtime_self_test() -> int:
    """Return a process exit code and persist diagnostics for windowed builds."""
    error_log = Path("ai-self-test-error.log")
    try:
        output = _ai_runtime_self_test()
    except Exception:
        error_log.write_text(
            ai_runtime_diagnostics() + "\n\n" + traceback.format_exc(),
            encoding="utf-8",
        )
        return 1
    error_log.unlink(missing_ok=True)
    Path("ai-self-test-diagnostics.log").write_text(
        ai_runtime_diagnostics()
        + f"\nSynthetic batch: OK\nAI Study Pack: OK\nAI_Knowledge: {output}\n",
        encoding="utf-8",
    )
    return 0


def run_ai_runtime_diagnostics() -> int:
    """Write and print a concise report while exercising the full batch path."""
    report = ai_runtime_diagnostics()
    error_log = Path("ai-self-test-error.log")
    try:
        output = _ai_runtime_self_test()
    except Exception:
        error_log.write_text(
            report + "\n\n" + traceback.format_exc(),
            encoding="utf-8",
        )
        return 1
    error_log.unlink(missing_ok=True)
    report += f"\nSynthetic batch: OK\nAI Study Pack: OK\nAI_Knowledge: {output}"
    Path("ai-self-test-diagnostics.log").write_text(report + "\n", encoding="utf-8")
    try:
        print(report)
    except (AttributeError, OSError, UnicodeError):
        pass
    return 0


def _load_ai_pipeline(script_path: Path) -> ModuleType:
    """Load the bundled local pipeline without spawning Python or the GUI EXE."""
    spec = importlib.util.spec_from_file_location("_bklms_ai_pipeline", script_path)
    if spec is None or spec.loader is None:
        raise AIPreparationError("Không thể tải thành phần chuẩn bị AI.")
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    try:
        spec.loader.exec_module(module)
    except Exception:
        sys.modules.pop(spec.name, None)
        raise
    if not callable(getattr(module, "run_preparation", None)):
        raise AIPreparationError("Thành phần chuẩn bị AI không hợp lệ.")
    return module


class AICoursePreparer:
    """Run the bundled local pipeline for exactly one downloaded course."""

    def __init__(
        self,
        *,
        script_path: Path | None = None,
        dependency_importer: Callable[[str], object] = importlib.import_module,
        pipeline_loader: Callable[[Path], ModuleType] = _load_ai_pipeline,
    ):
        self.script_path = script_path or default_ai_tool_path()
        self.dependency_importer = dependency_importer
        self.pipeline_loader = pipeline_loader
        self._pipeline: ModuleType | None = None

    def prepare(self, course_root: Path, output: Path | None = None) -> Path:
        missing = missing_ai_dependencies(self.dependency_importer)
        if missing:
            raise OptionalAIDependenciesError(
                "Thiếu thành phần AI cần thiết: " + ", ".join(missing) + "."
            )

        source_root = Path(course_root).expanduser().resolve()
        if not source_root.is_dir():
            raise AIPreparationError("Không tìm thấy thư mục course đã tải.")
        destination = Path(output or default_ai_output(source_root)).expanduser().resolve()
        expected_destination = default_ai_output(source_root).resolve()
        if destination != expected_destination:
            raise AIPreparationError("AI_Knowledge phải nằm trong thư mục course đã tải.")
        if not self.script_path.is_file():
            raise AIPreparationError("Không tìm thấy thành phần chuẩn bị AI trong ứng dụng.")

        try:
            pipeline = self._pipeline or self.pipeline_loader(self.script_path)
            self._pipeline = pipeline
            # The standalone tool prints Vietnamese CLI progress.  Redirect it
            # when embedded in the GUI so a Windows legacy console encoding can
            # never abort a local knowledge-base build.
            with redirect_stdout(io.StringIO()):
                pipeline.run_preparation(_pipeline_arguments(source_root, destination))
        except AIPreparationError:
            raise
        except Exception as exc:
            LOG.exception("AI preparation failed for %s", source_root)
            raise AIPreparationError("Không thể chuẩn bị course cho AI. Hãy thử lại sau.") from exc
        return destination


def _pipeline_arguments(course_root: Path, destination: Path) -> SimpleNamespace:
    """Keep GUI preparation local and deterministic: no transcription or cloud."""
    return SimpleNamespace(
        input=course_root,
        output=destination,
        include_references=False,
        transcribe=False,
        whisper_model="small",
        language=None,
        whisper_device="cpu",
        whisper_compute_type="int8",
        chunk_chars=4800,
        chunk_overlap=500,
        force=True,
    )


@dataclass(frozen=True)
class AICoursePreparationResult:
    course: Course
    output: Path | None = None
    error: str | None = None

    @property
    def succeeded(self) -> bool:
        return self.output is not None and self.error is None


@dataclass(frozen=True)
class AIBatchPreparationResult:
    results: list[AICoursePreparationResult]

    @property
    def succeeded(self) -> list[AICoursePreparationResult]:
        return [result for result in self.results if result.succeeded]

    @property
    def failed(self) -> list[AICoursePreparationResult]:
        return [result for result in self.results if not result.succeeded]


AIProgressCallback = Callable[[dict], None]
CourseRootResolver = Callable[[Course], Path]


class AIBatchPreparer:
    """Prepare independent per-course knowledge bases sequentially and safely."""

    def __init__(self, preparer_factory: Callable[[], AICoursePreparer] = AICoursePreparer):
        self.preparer_factory = preparer_factory

    def prepare_courses(
        self,
        courses: Iterable[Course],
        course_root_for: CourseRootResolver,
        progress_callback: AIProgressCallback | None = None,
    ) -> AIBatchPreparationResult:
        course_list = list(courses)
        if not course_list:
            return AIBatchPreparationResult([])

        preparer = self.preparer_factory()
        results: list[AICoursePreparationResult] = []
        total = len(course_list)
        for index, course in enumerate(course_list, start=1):
            self._emit(progress_callback, "ai_prepare_course_start", course=course, index=index, total=total)
            try:
                output = preparer.prepare(course_root_for(course))
                result = AICoursePreparationResult(course=course, output=output)
            except Exception as exc:
                LOG.exception("AI preparation failed for course %s", course.id)
                result = AICoursePreparationResult(
                    course=course,
                    error=_error_summary(exc),
                )
            results.append(result)
            self._emit(
                progress_callback,
                "ai_prepare_course_complete",
                result=result,
                index=index,
                total=total,
            )
        return AIBatchPreparationResult(results)

    @staticmethod
    def _emit(callback: AIProgressCallback | None, event: str, **payload) -> None:
        if callback is None:
            return
        try:
            callback({"event": event, **payload})
        except Exception:
            pass


def _error_summary(exc: Exception) -> str:
    message = f"{type(exc).__name__}: {exc}"
    if exc.__cause__ is not None:
        message += f" ({type(exc.__cause__).__name__}: {exc.__cause__})"
    return message
