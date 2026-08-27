#!/usr/bin/env python3
# -*- coding: utf-8 -*-
r"""
prepare_ai_course.py
====================
Biến thư mục tải thô từ BK-LMS Downloader thành knowledge base gọn, có cấu trúc,
phù hợp để feed cho AI tutor / RAG.

Không gọi LLM/API. Toàn bộ preprocessing chạy local.

Hỗ trợ:
- Input là thư mục hoặc .zip từ BK-LMS Downloader v2.
- content.txt / content.html (ưu tiên TXT, tránh duplicate HTML).
- PDF lecture: extract theo page.
- PPTX: extract theo slide.
- Video/audio: optional transcript bằng faster-whisper.
- .srt/.vtt: đưa transcript có sẵn vào corpus.
- .url: đưa vào links index, không feed làm knowledge text.
- Tự phân nhóm 00_course, Chapter 1..N, other.
- Tách reference books khỏi lecture corpus theo mặc định.
- Sinh Markdown chuẩn hóa + chunks JSONL để làm RAG.
- Sinh AI_TUTOR_CONTEXT.md, course_index.md, references_index.md,
  transcription_queue.md và manifest.

Ví dụ PowerShell:
    python prepare_ai_course.py `
      --input "C:\...\MMT\test_v2\Mạng máy tính (...)" `
      --output "C:\...\MMT\AI_Knowledge"

Transcript video:
    python prepare_ai_course.py ... --transcribe --whisper-model small

Muốn extract cả textbook/reference vào corpus:
    python prepare_ai_course.py ... --include-references
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import html
import json
import os
import re
import shutil
import sys
import tempfile
import textwrap
import unicodedata
import zipfile
from collections import Counter, defaultdict
from dataclasses import asdict, dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Iterable, Optional


_PROJECT_SRC = Path(__file__).resolve().parents[1] / "src"
if _PROJECT_SRC.is_dir() and str(_PROJECT_SRC) not in sys.path:
    sys.path.insert(0, str(_PROJECT_SRC))

from bklms_downloader.ai_study_pack import (
    create_chatgpt_study_pack,
    validate_ai_study_pack,
    write_study_navigation,
)


# -----------------------------------------------------------------------------
# Constants
# -----------------------------------------------------------------------------

TEXT_EXTS = {".txt", ".md"}
HTML_EXTS = {".html", ".htm"}
PDF_EXTS = {".pdf"}
PPT_EXTS = {".pptx"}
VIDEO_EXTS = {".mp4", ".mkv", ".webm", ".mov", ".avi"}
AUDIO_EXTS = {".mp3", ".m4a", ".wav", ".aac", ".flac", ".ogg"}
SUBTITLE_EXTS = {".srt", ".vtt"}
URL_EXTS = {".url"}
JSON_EXTS = {".json"}

SKIP_DIR_NAMES = {
    "AI_Knowledge",
    "__MACOSX",
    ".git",
    ".idea",
    ".vscode",
    "node_modules",
    "duylms_forum_debug",
}
STUDY_PACK_SUFFIX = " - AI Study Pack.zip"

REFERENCE_HINTS = (
    "textbook",
    "books and references",
    "book and reference",
    "references",
    "cisco ccie",
    "routing tcp",
    "prentice hall",
    "top-down approach",
    "top down approach",
)

COURSE_META_HINTS = (
    "course syllabus",
    "general information",
    "course outline",
    "grading",
    "final exam",
    "thi cuối kỳ",
    "thông tin môn học",
    "overview",
    "description",
    "books and references",
    "announcements",
    "chung",
)

# Dùng để xếp ưu tiên nguồn khi AI tutor retrieval.
SOURCE_PRIORITY = {
    "lms_page": 1,
    "lms_text": 1,
    "slide": 2,
    "lecture_pdf": 2,
    "subtitle": 3,
    "video_transcript": 3,
    "reference_pdf": 4,
    "reference_text": 4,
    "other": 5,
}


# -----------------------------------------------------------------------------
# Data model
# -----------------------------------------------------------------------------

@dataclass
class DocumentRecord:
    source_id: str
    title: str
    source_path: str
    source_type: str
    priority: int
    group: str
    chapter: Optional[int]
    output_path: Optional[str]
    status: str
    chars: int = 0
    words: int = 0
    units: int = 0  # pages/slides/segments/etc.
    note: str = ""
    chapters: list[int] = field(default_factory=list)
    order: int = 0
    source_copy_path: Optional[str] = None


@dataclass
class ChunkRecord:
    chunk_id: str
    source_id: str
    title: str
    source_path: str
    source_type: str
    priority: int
    group: str
    chapter: Optional[int]
    locator: str
    text: str
    chunk_path: str
    chapters: list[int] = field(default_factory=list)


# -----------------------------------------------------------------------------
# Generic helpers
# -----------------------------------------------------------------------------

def clean_text(value: str) -> str:
    value = (value or "").replace("\xa0", " ")
    value = value.replace("\r\n", "\n").replace("\r", "\n")
    # Giữ paragraph breaks nhưng dọn whitespace nội dòng.
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in value.split("\n")]
    out: list[str] = []
    blank = False
    for line in lines:
        if line:
            out.append(line)
            blank = False
        elif out and not blank:
            out.append("")
            blank = True
    return "\n".join(out).strip()


def normalized_for_compare(value: str) -> str:
    value = unicodedata.normalize("NFKC", value or "").lower()
    value = re.sub(r"\s+", " ", value)
    return value.strip()


def safe_name(value: str, max_len: int = 120) -> str:
    value = clean_text(value)
    value = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", value)
    value = re.sub(r"\s+", " ", value).strip(" .")
    if not value:
        value = "untitled"
    if len(value) > max_len:
        value = value[:max_len].rstrip(" .")
    return value


def slug(value: str, max_len: int = 80) -> str:
    value = unicodedata.normalize("NFKD", value or "")
    value = "".join(ch for ch in value if not unicodedata.combining(ch))
    value = value.lower()
    value = re.sub(r"[^a-z0-9]+", "_", value).strip("_")
    return (value or "document")[:max_len]


def short_hash(value: str, length: int = 10) -> str:
    return hashlib.sha1(value.encode("utf-8", errors="ignore")).hexdigest()[:length]


def rel(path: Path, root: Path) -> str:
    try:
        return path.relative_to(root).as_posix()
    except Exception:
        return path.as_posix()


def write_text(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text.rstrip() + "\n", encoding="utf-8")


def word_count(text: str) -> int:
    return len(re.findall(r"\S+", text or ""))


def yaml_escape(value: str) -> str:
    return json.dumps(value, ensure_ascii=False)


def path_haystack(path: Path) -> str:
    return normalized_for_compare(" / ".join(path.parts))


def is_reference_path(path: Path) -> bool:
    h = path_haystack(path)
    return any(hint in h for hint in REFERENCE_HINTS)


def chapter_numbers_from_path(path: Path) -> list[int]:
    """Detect explicit chapter numbers/ranges without mistaking order prefixes."""
    text = " / ".join(path.parts)
    range_pattern = re.compile(
        r"(?i)(?<![a-z])(?:chapter|ch)\s*0*(\d{1,2})\s*"
        r"(?:[_\-–—&]|\band\b|\bto\b)\s*"
        r"(?:(?:chapter|ch)\s*)?0*(\d{1,2})"
    )
    match = range_pattern.search(text)
    if match:
        start, end = sorted((int(match.group(1)), int(match.group(2))))
        if 0 < start < 100 and 0 < end < 100:
            return list(range(start, end + 1))

    single_pattern = re.compile(r"(?i)(?<![a-z])(?:chapter|ch)\s*0*(\d{1,2})")
    matches = [int(match.group(1)) for match in single_pattern.finditer(text)]
    matches = [number for number in matches if 0 < number < 100]
    return [matches[-1]] if matches else []


def chapter_numbers_from_group(group: str, chapter: Optional[int] = None) -> list[int]:
    values = [int(value) for value in re.findall(r"\d+", group or "")]
    if values:
        return values
    return [chapter] if chapter is not None else []


def chapter_from_path(path: Path) -> Optional[int]:
    numbers = chapter_numbers_from_path(path)
    return numbers[0] if numbers else None


def classify_group(path: Path) -> tuple[str, Optional[int]]:
    chapters = chapter_numbers_from_path(path)
    if chapters:
        return "chapter_" + "_".join(f"{chapter:02d}" for chapter in chapters), chapters[0]

    h = path_haystack(path)
    if any(hint in h for hint in COURSE_META_HINTS):
        return "00_course", None

    if is_reference_path(path):
        return "references", None

    return "other", None


def source_id_for(path: Path, root: Path, source_type: str) -> str:
    key = f"{rel(path, root)}|{source_type}"
    return f"src_{short_hash(key, 12)}"


def source_order_hint(source_path: str) -> int:
    match = re.match(r"\s*0*(\d{1,3})(?:[_ .-]|$)", Path(source_path).name)
    return int(match.group(1)) if match else 9999


def choose_group_dir(kb_docs_dir: Path, group: str) -> Path:
    if group == "00_course":
        return kb_docs_dir / "00_course"
    if group.startswith("chapter_"):
        return kb_docs_dir / "chapters" / group
    if group == "references":
        return kb_docs_dir / "references"
    return kb_docs_dir / "other"


def infer_title(path: Path) -> str:
    if path.name.lower() in {"content.txt", "content.html", "content.htm"}:
        return clean_text(path.parent.name) or path.stem
    return clean_text(path.stem.replace("_", " ")) or path.name


# -----------------------------------------------------------------------------
# Input discovery / unzip
# -----------------------------------------------------------------------------

def find_course_root(root: Path) -> Path:
    """Nếu ZIP có wrapper test_v2/<course>/..., chọn root hữu ích nhất."""
    current = root
    for _ in range(4):
        children = [p for p in current.iterdir() if p.name not in SKIP_DIR_NAMES]
        dirs = [p for p in children if p.is_dir()]
        files = [p for p in children if p.is_file()]

        # Một wrapper chỉ có đúng 1 folder và hầu như không có file.
        if len(dirs) == 1 and not files:
            current = dirs[0]
            continue
        break
    return current


def prepare_input(input_path: Path, work_dir: Path) -> tuple[Path, Optional[Path]]:
    input_path = input_path.expanduser().resolve()
    if not input_path.exists():
        raise FileNotFoundError(f"Không tồn tại input: {input_path}")

    if input_path.is_dir():
        return find_course_root(input_path), None

    if input_path.suffix.lower() == ".zip":
        extracted = work_dir / "unzipped"
        extracted.mkdir(parents=True, exist_ok=True)
        print(f"[1/6] Giải nén: {input_path.name}")
        with zipfile.ZipFile(input_path, "r") as zf:
            zf.extractall(extracted)
        return find_course_root(extracted), extracted

    raise ValueError("--input phải là thư mục hoặc file .zip")


def iter_source_files(root: Path) -> list[Path]:
    results: list[Path] = []
    for path in root.rglob("*"):
        if not path.is_file():
            continue
        if any(part in SKIP_DIR_NAMES for part in path.parts):
            continue
        if path.name.endswith(STUDY_PACK_SUFFIX):
            continue
        # Ignore downloader metadata as teaching content; still index separately later.
        results.append(path)
    return sorted(results, key=lambda p: p.as_posix().lower())


# -----------------------------------------------------------------------------
# Extractors
# -----------------------------------------------------------------------------

def read_plain_text(path: Path) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp1258", "cp1252", "latin-1"):
        try:
            return clean_text(path.read_text(encoding=enc))
        except UnicodeDecodeError:
            continue
    return clean_text(path.read_text(encoding="utf-8", errors="replace"))


def html_to_markdown(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:
        raise RuntimeError("Thiếu beautifulsoup4. Chạy: pip install beautifulsoup4") from exc

    soup = BeautifulSoup(raw, "html.parser")
    for node in soup.select("script, style, noscript, nav, header, footer, button, form"):
        node.decompose()

    # Ưu tiên markdownify nếu có để giữ headings/lists/links.
    try:
        from markdownify import markdownify as md
        text = md(str(soup), heading_style="ATX")
        return clean_text(text)
    except ImportError:
        return clean_text(soup.get_text("\n", strip=True))


def extract_pdf(path: Path) -> tuple[str, list[tuple[str, str]]]:
    """Return (full markdown body, [(locator, unit_text), ...])."""
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("Thiếu pypdf. Chạy: pip install pypdf") from exc

    reader = PdfReader(path)
    units: list[tuple[str, str]] = []
    pieces: list[str] = []

    for idx, page in enumerate(reader.pages, start=1):
        text = clean_text(page.extract_text() or "")
        if not text:
            continue
        locator = f"page {idx}"
        units.append((locator, text))
        pieces.append(f"## Page {idx}\n\n{text}")

    if not pieces:
        body = (
            "_Không extract được text từ PDF này. Có thể PDF là scan/image. "
            "Tool không OCR mặc định._"
        )
    else:
        body = "\n\n".join(pieces)

    return body, units


def extract_pptx(path: Path) -> tuple[str, list[tuple[str, str]]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("Thiếu python-pptx. Chạy: pip install python-pptx") from exc

    prs = Presentation(path)
    units: list[tuple[str, str]] = []
    pieces: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = clean_text(getattr(shape, "text", ""))
                if t:
                    texts.append(t)

            # Table text.
            if getattr(shape, "has_table", False):
                try:
                    for row in shape.table.rows:
                        row_text = " | ".join(clean_text(c.text) for c in row.cells)
                        if row_text.strip(" |"):  # pragma: no branch
                            texts.append(row_text)
                except Exception:
                    pass

        # Notes nếu python-pptx version hỗ trợ.
        try:
            notes = slide.notes_slide
            note_texts: list[str] = []
            for shape in notes.shapes:
                if hasattr(shape, "text"):
                    t = clean_text(getattr(shape, "text", ""))
                    if t and t.lower() not in {"slide image", "slide number"}:
                        note_texts.append(t)
            if note_texts:
                texts.append("[Speaker notes]\n" + "\n".join(note_texts))
        except Exception:
            pass

        text = clean_text("\n".join(texts))
        if not text:
            continue

        locator = f"slide {idx}"
        units.append((locator, text))
        pieces.append(f"## Slide {idx}\n\n{text}")

    return "\n\n".join(pieces), units


def parse_subtitle(path: Path) -> tuple[str, list[tuple[str, str]]]:
    raw = read_plain_text(path)
    # VTT header.
    raw = re.sub(r"^WEBVTT.*?\n+", "", raw, flags=re.I | re.S)
    blocks = re.split(r"\n\s*\n", raw)
    units: list[tuple[str, str]] = []
    pieces: list[str] = []

    time_re = re.compile(
        r"(?P<start>\d{1,2}:\d{2}(?::\d{2})?[,.]\d{3})\s*-->\s*"
        r"(?P<end>\d{1,2}:\d{2}(?::\d{2})?[,.]\d{3})"
    )

    for block in blocks:
        lines = [l.strip() for l in block.splitlines() if l.strip()]
        if not lines:
            continue
        time_idx = next((i for i, l in enumerate(lines) if "-->" in l), None)
        if time_idx is None:
            continue
        m = time_re.search(lines[time_idx])
        locator = m.group("start") if m else lines[time_idx].split("-->")[0].strip()
        text = clean_text(" ".join(lines[time_idx + 1 :]))
        if not text:
            continue
        units.append((locator, text))
        pieces.append(f"[{locator}] {text}")

    return "\n\n".join(pieces), units


def transcribe_media(
    path: Path,
    model_name: str,
    language: Optional[str],
    device: str,
    compute_type: str,
    model_cache: dict,
) -> tuple[str, list[tuple[str, str]]]:
    try:
        from faster_whisper import WhisperModel
    except ImportError as exc:
        raise RuntimeError(
            "Thiếu faster-whisper. Chạy: pip install faster-whisper"
        ) from exc

    cache_key = (model_name, device, compute_type)
    if cache_key not in model_cache:
        print(f"    Loading Whisper model '{model_name}' ({device}, {compute_type})...")
        model_cache[cache_key] = WhisperModel(
            model_name,
            device=device,
            compute_type=compute_type,
        )

    model = model_cache[cache_key]
    segments, info = model.transcribe(
        str(path),
        language=language,
        vad_filter=True,
        beam_size=5,
    )

    units: list[tuple[str, str]] = []
    pieces: list[str] = []

    for seg in segments:
        text = clean_text(seg.text)
        if not text:
            continue
        locator = seconds_to_timestamp(float(seg.start))
        units.append((locator, text))
        pieces.append(f"[{locator}] {text}")

    header = ""
    try:
        header = f"Detected language: {info.language} (p={info.language_probability:.3f})\n\n"
    except Exception:
        pass

    return header + "\n".join(pieces), units


def seconds_to_timestamp(sec: float) -> str:
    sec = max(0, int(sec))
    h = sec // 3600
    m = (sec % 3600) // 60
    s = sec % 60
    return f"{h:02d}:{m:02d}:{s:02d}"


def parse_internet_shortcut(path: Path) -> Optional[str]:
    text = read_plain_text(path)
    m = re.search(r"(?im)^URL=(.+)$", text)
    return m.group(1).strip() if m else None


# -----------------------------------------------------------------------------
# Chunking
# -----------------------------------------------------------------------------

def split_unit_text(text: str, max_chars: int, overlap_chars: int) -> list[str]:
    text = clean_text(text)
    if not text:
        return []
    if len(text) <= max_chars:
        return [text]

    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    chunks: list[str] = []
    current = ""

    for p in paragraphs:
        if len(p) > max_chars:
            # Flush current first.
            if current:
                chunks.append(current.strip())
                current = ""
            start = 0
            while start < len(p):
                end = min(len(p), start + max_chars)
                piece = p[start:end].strip()
                if piece:
                    chunks.append(piece)
                if end >= len(p):
                    break
                start = max(0, end - overlap_chars)
            continue

        candidate = p if not current else current + "\n\n" + p
        if len(candidate) <= max_chars:
            current = candidate
        else:
            chunks.append(current.strip())
            tail = current[-overlap_chars:].strip() if overlap_chars > 0 else ""
            current = (tail + "\n\n" + p).strip() if tail else p

    if current:
        chunks.append(current.strip())

    return [c for c in chunks if c]


def make_chunks(
    source_id: str,
    title: str,
    source_path: str,
    source_type: str,
    priority: int,
    group: str,
    chapter: Optional[int],
    chapters: list[int],
    units: list[tuple[str, str]],
    fallback_text: str,
    chunk_dir: Path,
    max_chars: int,
    overlap_chars: int,
) -> list[ChunkRecord]:
    if not units and fallback_text:
        units = [("document", fallback_text)]

    results: list[ChunkRecord] = []
    ordinal = 0

    for locator, unit_text in units:
        for piece in split_unit_text(unit_text, max_chars, overlap_chars):
            ordinal += 1
            chunk_id = f"{source_id}_c{ordinal:04d}"
            group_dir = chunk_dir / group
            if group.startswith("chapter_"):
                group_dir = chunk_dir / "chapters" / group
            chunk_path = group_dir / f"{chunk_id}.md"

            md = (
                "---\n"
                f"chunk_id: {yaml_escape(chunk_id)}\n"
                f"source_id: {yaml_escape(source_id)}\n"
                f"title: {yaml_escape(title)}\n"
                f"source_path: {yaml_escape(source_path)}\n"
                f"source_type: {yaml_escape(source_type)}\n"
                f"priority: {priority}\n"
                f"group: {yaml_escape(group)}\n"
                f"chapter: {json.dumps(chapter)}\n"
                f"chapters: {json.dumps(chapters)}\n"
                f"locator: {yaml_escape(locator)}\n"
                "---\n\n"
                f"# {title}\n\n"
                f"**Location:** {locator}\n\n"
                f"{piece}\n"
            )
            write_text(chunk_path, md)

            results.append(
                ChunkRecord(
                    chunk_id=chunk_id,
                    source_id=source_id,
                    title=title,
                    source_path=source_path,
                    source_type=source_type,
                    priority=priority,
                    group=group,
                    chapter=chapter,
                    chapters=list(chapters),
                    locator=locator,
                    text=piece,
                    chunk_path=chunk_path.as_posix(),
                )
            )

    return results


# -----------------------------------------------------------------------------
# Document writing
# -----------------------------------------------------------------------------

def write_document_markdown(
    output_path: Path,
    source_id: str,
    title: str,
    source_path: str,
    source_type: str,
    priority: int,
    group: str,
    chapter: Optional[int],
    chapters: list[int],
    body: str,
) -> None:
    front = (
        "---\n"
        f"source_id: {yaml_escape(source_id)}\n"
        f"title: {yaml_escape(title)}\n"
        f"source_path: {yaml_escape(source_path)}\n"
        f"source_type: {yaml_escape(source_type)}\n"
        f"priority: {priority}\n"
        f"group: {yaml_escape(group)}\n"
        f"chapter: {json.dumps(chapter)}\n"
        f"chapters: {json.dumps(chapters)}\n"
        "---\n\n"
    )
    write_text(output_path, front + f"# {title}\n\n" + body)


# -----------------------------------------------------------------------------
# Main processor
# -----------------------------------------------------------------------------

def is_media_activity_stub(path: Path, text: str) -> bool:
    """Skip tiny Moodle page text whose real content is a local media asset."""
    if word_count(text) >= 80:
        return False

    parent = path.parent
    # Chỉ nhìn tên activity gần nhất; không nhìn toàn ancestor path vì course phụ
    # có thể có chữ "Video" trong tên và làm mọi page bị nhận nhầm.
    media_named = bool(re.search(r"(?i)(^|[ _-])video([ _-]|$)", parent.name))
    assets = parent / "assets"
    has_media = False
    if assets.exists():
        try:
            has_media = any(
                x.is_file() and x.suffix.lower() in (VIDEO_EXTS | AUDIO_EXTS)
                for x in assets.iterdir()
            )
        except Exception:
            pass
    return media_named or has_media


class CoursePreparer:
    def __init__(self, args, source_root: Path, output_root: Path):
        self.args = args
        self.source_root = source_root
        self.output_root = output_root

        self.docs_dir = output_root / "documents"
        self.chunks_dir = output_root / "chunks"
        self.meta_dir = output_root / "meta"
        self.sources_dir = output_root / "sources"

        self.records: list[DocumentRecord] = []
        self.chunks: list[ChunkRecord] = []
        self.links: list[dict] = []
        self.reference_paths: list[dict] = []
        self.transcription_queue: list[dict] = []
        self.errors: list[dict] = []
        self.visual_sources: list[dict] = []
        self.seen_content_hashes: dict[str, str] = {}
        self.whisper_model_cache: dict = {}

        self.course_name = source_root.name

    def register_record(self, record: DocumentRecord) -> None:
        if not record.chapters:
            record.chapters = chapter_numbers_from_group(record.group, record.chapter)
        if not record.order:
            record.order = source_order_hint(record.source_path)
        self.records.append(record)

    def retain_visual_source(self, path: Path, source_id: str, source_type: str) -> Optional[str]:
        """Keep original lecturer visuals available without treating text as visual truth."""
        if source_type not in {"lecture_pdf", "slide"}:
            return None
        destination = self.sources_dir / f"{source_id}__{safe_name(path.name, 150)}"
        destination.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(path, destination)
        copied_path = rel(destination, self.output_root)
        self.visual_sources.append(
            {
                "source_id": source_id,
                "source_path": rel(path, self.source_root),
                "copied_path": copied_path,
                "source_type": source_type,
                "note": "Original lecturer source retained for diagrams and layout-dependent meaning.",
            }
        )
        return copied_path

    def is_duplicate_text(self, text: str, source_path: str) -> Optional[str]:
        normalized = normalized_for_compare(text)
        if len(normalized) < 80:
            return None
        digest = hashlib.sha1(normalized.encode("utf-8")).hexdigest()
        first = self.seen_content_hashes.get(digest)
        if first:
            return first
        self.seen_content_hashes[digest] = source_path
        return None

    def output_doc_path(self, source_id: str, title: str, group: str) -> Path:
        d = choose_group_dir(self.docs_dir, group)
        return d / f"{slug(title)}__{source_id}.md"

    def add_text_document(
        self,
        path: Path,
        source_type: str,
        body: str,
        units: Optional[list[tuple[str, str]]] = None,
        note: str = "",
    ) -> None:
        body = clean_text(body)
        source_path = rel(path, self.source_root)
        title = infer_title(path)
        group, chapter = classify_group(path)
        chapters = chapter_numbers_from_group(group, chapter)
        priority = SOURCE_PRIORITY.get(source_type, 5)
        source_id = source_id_for(path, self.source_root, source_type)

        if not body:
            self.register_record(
                DocumentRecord(
                    source_id=source_id,
                    title=title,
                    source_path=source_path,
                    source_type=source_type,
                    priority=priority,
                    group=group,
                    chapter=chapter,
                    output_path=None,
                    status="empty",
                    note=note,
                )
            )
            return

        duplicate_of = self.is_duplicate_text(body, source_path)
        if duplicate_of:
            self.register_record(
                DocumentRecord(
                    source_id=source_id,
                    title=title,
                    source_path=source_path,
                    source_type=source_type,
                    priority=priority,
                    group=group,
                    chapter=chapter,
                    output_path=None,
                    status="duplicate",
                    chars=len(body),
                    words=word_count(body),
                    units=len(units or []),
                    note=f"Duplicate text of {duplicate_of}",
                )
            )
            return

        out_path = self.output_doc_path(source_id, title, group)
        write_document_markdown(
            out_path,
            source_id,
            title,
            source_path,
            source_type,
            priority,
            group,
            chapter,
            chapters,
            body,
        )

        chunk_records = make_chunks(
            source_id=source_id,
            title=title,
            source_path=source_path,
            source_type=source_type,
            priority=priority,
            group=group,
            chapter=chapter,
            chapters=chapters,
            units=units or [],
            fallback_text=body,
            chunk_dir=self.chunks_dir,
            max_chars=self.args.chunk_chars,
            overlap_chars=self.args.chunk_overlap,
        )

        self.chunks.extend(chunk_records)
        self.register_record(
            DocumentRecord(
                source_id=source_id,
                title=title,
                source_path=source_path,
                source_type=source_type,
                priority=priority,
                group=group,
                chapter=chapter,
                output_path=rel(out_path, self.output_root),
                status="ready",
                chars=len(body),
                words=word_count(body),
                units=len(units or []),
                note=note,
                chapters=chapters,
                source_copy_path=self.retain_visual_source(path, source_id, source_type),
            )
        )

    def process_content_txt(self, path: Path) -> None:
        text = read_plain_text(path)

        # Moodle video activities often produce a tiny content.txt containing only
        # the activity label/embed wrapper. Feeding that as Priority 1 harms retrieval;
        # the actual video belongs in the transcript layer (Priority 3).
        if is_media_activity_stub(path, text):
            source_path = rel(path, self.source_root)
            sid = source_id_for(path, self.source_root, "lms_page")
            group, chapter = classify_group(path)
            self.register_record(
                DocumentRecord(
                    source_id=sid,
                    title=infer_title(path),
                    source_path=source_path,
                    source_type="lms_page",
                    priority=1,
                    group=group,
                    chapter=chapter,
                    output_path=None,
                    status="skipped_media_stub",
                    chars=len(text),
                    words=word_count(text),
                    note="Tiny activity wrapper; use the media transcript instead",
                )
            )
            return

        self.add_text_document(path, "lms_page", text)

    def process_html(self, path: Path) -> None:
        # Nếu cùng thư mục đã có content.txt thì TXT là canonical; bỏ HTML duplicate.
        sibling_txt = path.with_name("content.txt")
        if path.name.lower().startswith("content.") and sibling_txt.exists():
            source_path = rel(path, self.source_root)
            sid = source_id_for(path, self.source_root, "lms_page")
            group, chapter = classify_group(path)
            self.register_record(
                DocumentRecord(
                    source_id=sid,
                    title=infer_title(path),
                    source_path=source_path,
                    source_type="lms_page",
                    priority=1,
                    group=group,
                    chapter=chapter,
                    output_path=None,
                    status="skipped_html_duplicate",
                    note="content.txt exists in same folder",
                )
            )
            return
        self.add_text_document(path, "lms_page", html_to_markdown(path))

    def process_pdf(self, path: Path) -> None:
        ref = is_reference_path(path)
        source_type = "reference_pdf" if ref else "lecture_pdf"

        if ref and not self.args.include_references:
            self.reference_paths.append(
                {
                    "title": infer_title(path),
                    "source_path": rel(path, self.source_root),
                    "size_mb": round(path.stat().st_size / (1024 * 1024), 2),
                    "status": "indexed_only",
                }
            )
            group, chapter = classify_group(path)
            sid = source_id_for(path, self.source_root, source_type)
            self.register_record(
                DocumentRecord(
                    source_id=sid,
                    title=infer_title(path),
                    source_path=rel(path, self.source_root),
                    source_type=source_type,
                    priority=4,
                    group=group,
                    chapter=chapter,
                    output_path=None,
                    status="reference_indexed_only",
                    note="Use --include-references to extract this PDF",
                )
            )
            return

        body, units = extract_pdf(path)
        self.add_text_document(path, source_type, body, units)

    def process_pptx(self, path: Path) -> None:
        body, units = extract_pptx(path)
        self.add_text_document(path, "slide", body, units)

    def process_subtitle(self, path: Path) -> None:
        body, units = parse_subtitle(path)
        self.add_text_document(path, "subtitle", body, units)

    def process_media(self, path: Path) -> None:
        if not self.args.transcribe:
            group, chapter = classify_group(path)
            source_path = rel(path, self.source_root)
            source_id = source_id_for(path, self.source_root, "media")
            self.transcription_queue.append(
                {
                    "title": infer_title(path),
                    "source_path": source_path,
                    "size_mb": round(path.stat().st_size / (1024 * 1024), 2),
                }
            )
            self.register_record(
                DocumentRecord(
                    source_id=source_id,
                    title=infer_title(path),
                    source_path=source_path,
                    source_type="media",
                    priority=3,
                    group=group,
                    chapter=chapter,
                    output_path=None,
                    status="media_pending",
                    note="Media was not transcribed; see transcription_queue.md",
                )
            )
            return

        body, units = transcribe_media(
            path=path,
            model_name=self.args.whisper_model,
            language=self.args.language,
            device=self.args.whisper_device,
            compute_type=self.args.whisper_compute_type,
            model_cache=self.whisper_model_cache,
        )
        self.add_text_document(path, "video_transcript", body, units)

    def process_url(self, path: Path) -> None:
        group, chapter = classify_group(path)
        source_path = rel(path, self.source_root)
        target = parse_internet_shortcut(path)
        self.links.append(
            {
                "title": infer_title(path),
                "source_path": source_path,
                "url": target,
                "group": group,
            }
        )
        self.register_record(
            DocumentRecord(
                source_id=source_id_for(path, self.source_root, "url"),
                title=infer_title(path),
                source_path=source_path,
                source_type="url",
                priority=5,
                group=group,
                chapter=chapter,
                output_path=None,
                status="link_only",
                note=f"Unresolved external content: {target or 'URL not parsed'}",
            )
        )

    def process_json_metadata(self, path: Path) -> None:
        if path.name not in {"_course_structure.json", "_stats.json", "_download_manifest.json"}:
            self.process_unhandled(path)
            return
        try:
            data = json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            return
        out = self.meta_dir / "raw_downloader_metadata" / safe_name(path.name)
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
        group, chapter = classify_group(path)
        self.register_record(
            DocumentRecord(
                source_id=source_id_for(path, self.source_root, "metadata"),
                title=infer_title(path),
                source_path=rel(path, self.source_root),
                source_type="metadata",
                priority=5,
                group=group,
                chapter=chapter,
                output_path=rel(out, self.output_root),
                status="metadata_copied",
                note="Downloader metadata retained for traceability, not teaching content.",
            )
        )

    def process_unhandled(self, path: Path) -> None:
        group, chapter = classify_group(path)
        self.register_record(
            DocumentRecord(
                source_id=source_id_for(path, self.source_root, "binary"),
                title=infer_title(path),
                source_path=rel(path, self.source_root),
                source_type="binary",
                priority=5,
                group=group,
                chapter=chapter,
                output_path=None,
                status="skipped_unsupported",
                note=f"Unsupported source type: {path.suffix.lower() or 'no extension'}",
            )
        )

    def process(self) -> None:
        files = iter_source_files(self.source_root)
        print(f"[2/6] Source root: {self.source_root}")
        print(f"      Phát hiện {len(files)} file")

        # Canonical order: content.txt trước HTML để dedup chính xác hơn.
        def order_key(p: Path):
            ext = p.suffix.lower()
            if p.name.lower() == "content.txt":
                rank = 0
            elif ext in TEXT_EXTS:
                rank = 1
            elif ext in HTML_EXTS:
                rank = 2
            elif ext in PPT_EXTS:
                rank = 3
            elif ext in PDF_EXTS:
                rank = 4
            elif ext in SUBTITLE_EXTS:
                rank = 5
            elif ext in VIDEO_EXTS | AUDIO_EXTS:
                rank = 6
            elif ext in URL_EXTS:
                rank = 7
            else:
                rank = 9
            return (rank, p.as_posix().lower())

        for idx, path in enumerate(sorted(files, key=order_key), start=1):
            ext = path.suffix.lower()
            source_path = rel(path, self.source_root)

            # Downloader metadata copy.
            if ext in JSON_EXTS:
                self.process_json_metadata(path)
                continue

            # Ignore generated helper text from old downloader if obviously empty/noisy.
            try:
                print(f"      [{idx:03d}/{len(files):03d}] {source_path}")
                if ext in TEXT_EXTS:
                    self.process_content_txt(path)
                elif ext in HTML_EXTS:
                    self.process_html(path)
                elif ext in PDF_EXTS:
                    self.process_pdf(path)
                elif ext in PPT_EXTS:
                    self.process_pptx(path)
                elif ext in SUBTITLE_EXTS:
                    self.process_subtitle(path)
                elif ext in VIDEO_EXTS | AUDIO_EXTS:
                    self.process_media(path)
                elif ext in URL_EXTS:
                    self.process_url(path)
                else:
                    self.process_unhandled(path)
            except Exception as exc:
                print(f"        [WARN] {exc}")
                self.errors.append({"source_path": source_path, "error": str(exc)})
                group, chapter = classify_group(path)
                sid = source_id_for(path, self.source_root, "other")
                self.register_record(
                    DocumentRecord(
                        source_id=sid,
                        title=infer_title(path),
                        source_path=source_path,
                        source_type="other",
                        priority=5,
                        group=group,
                        chapter=chapter,
                        output_path=None,
                        status="error",
                        note=str(exc),
                    )
                )

        print("[3/6] Sinh manifests/chunks...")
        self.write_manifests()
        print("[4/6] Sinh course_index.md...")
        self.write_course_index()
        print("[5/6] Sinh AI_TUTOR_CONTEXT.md...")
        self.write_ai_tutor_context()
        print("[6/6] Sinh reports phụ...")
        self.write_reports()
        self.write_study_pack_outputs()

    def write_manifests(self) -> None:
        self.meta_dir.mkdir(parents=True, exist_ok=True)

        with (self.meta_dir / "documents.jsonl").open("w", encoding="utf-8") as f:
            for r in self.records:
                f.write(json.dumps(asdict(r), ensure_ascii=False) + "\n")

        # JSONL ready for embedding/import. chunk_path stored relative to KB root.
        with (self.meta_dir / "corpus.jsonl").open("w", encoding="utf-8") as f:
            for c in self.chunks:
                item = asdict(c)
                try:
                    item["chunk_path"] = rel(Path(c.chunk_path), self.output_root)
                except Exception:
                    pass
                f.write(json.dumps(item, ensure_ascii=False) + "\n")

        # CSV convenient for manual review.
        with (self.meta_dir / "documents.csv").open("w", encoding="utf-8-sig", newline="") as f:
            fieldnames = list(DocumentRecord.__dataclass_fields__.keys())
            w = csv.DictWriter(f, fieldnames=fieldnames)
            w.writeheader()
            for r in self.records:
                w.writerow(asdict(r))

        stats = {
            "course_name": self.course_name,
            "source_root": ".",
            "source_paths_relative": True,
            "generated_at": datetime.now().isoformat(timespec="seconds"),
            "documents_total": len(self.records),
            "documents_ready": sum(r.status == "ready" for r in self.records),
            "chunks": len(self.chunks),
            "words_ready": sum(r.words for r in self.records if r.status == "ready"),
            "references_indexed_only": len(self.reference_paths),
            "media_pending_transcription": len(self.transcription_queue),
            "links": len(self.links),
            "errors": len(self.errors),
            "source_types": dict(Counter(r.source_type for r in self.records)),
            "statuses": dict(Counter(r.status for r in self.records)),
            "chapter_groups": sorted({r.group for r in self.records if r.group.startswith("chapter_")}),
        }
        (self.meta_dir / "stats.json").write_text(
            json.dumps(stats, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

    def write_course_index(self) -> None:
        ready = [r for r in self.records if r.status == "ready"]
        groups: dict[str, list[DocumentRecord]] = defaultdict(list)
        for r in ready:
            groups[r.group].append(r)

        def group_sort_key(name: str):
            if name == "00_course":
                return (0, 0)
            m = re.match(r"chapter_(\d+)", name)
            if m:
                return (1, int(m.group(1)))
            if name == "references":
                return (2, 0)
            return (3, 0)

        lines = [
            f"# Course Index — {self.course_name}",
            "",
            "Knowledge base generated from the raw LMS download.",
            "",
            "## Source priority",
            "",
            "1. LMS pages / lecturer-authored course text",
            "2. Lecture slides / lecture PDFs",
            "3. Video transcripts / subtitles",
            "4. Textbooks and external references",
            "",
        ]

        for group_name in sorted(groups, key=group_sort_key):
            label = (
                "Course information"
                if group_name == "00_course"
                else group_name.replace("_", " ").title()
            )
            lines.extend([f"## {label}", ""])
            for r in sorted(groups[group_name], key=lambda x: (x.priority, x.title.lower())):
                lines.append(
                    f"- **P{r.priority} · {r.source_type}** — {r.title} "
                    f"(`{r.source_id}`; `{r.source_path}`)"
                )
            lines.append("")

        write_text(self.output_root / "course_index.md", "\n".join(lines))

    def write_ai_tutor_context(self) -> None:
        chapter_sets = sorted(
            {tuple(record.chapters) for record in self.records if record.chapters},
            key=lambda values: values,
        )
        chapter_text = ", ".join(
            str(values[0]) if len(values) == 1 else "–".join(map(str, values))
            for values in chapter_sets
        ) or "not detected"
        detected = {chapter for values in chapter_sets for chapter in values}
        gaps = [chapter for chapter in range(min(detected), max(detected) + 1) if chapter not in detected] if detected else []

        text = f"""# AI Tutor Context — {self.course_name}

## Purpose

This folder is an AI-ready knowledge base produced from the course's LMS materials. Use it as the primary evidence when teaching, reviewing, generating quizzes, or answering course-specific questions.

## Evidence priority

When multiple sources overlap or disagree, prefer them in this order:

1. **Priority 1 — LMS pages / lecturer-authored course text.** These best reflect the instructor's framing, course rules, grading, exam scope, and terminology.
2. **Priority 2 — Lecture slides and lecture PDFs.** Use these for the concepts actually emphasized in lectures.
3. **Priority 3 — Video transcripts / subtitles.** Use these to recover oral explanations and examples from lectures.
4. **Priority 4 — Textbooks / external references.** Use these for deeper explanation or missing detail, but do not let broader textbook coverage silently override the course scope.

## Teaching rules

- Ground course-specific claims in the supplied sources; do not silently replace the instructor's framing with generic knowledge.
- When a claim comes from a chunk, cite **source_id + locator** (for example `src_xxx, slide 12`, `src_xxx, page 45`, or `src_xxx, 00:12:31`).
- If the supplied materials do not support an answer, say that clearly before adding general knowledge.
- Distinguish **course requirement/exam scope** from **reference enrichment**.
- Prefer explaining from the chapter's priority 1–3 sources first; retrieve priority 4 only when useful.
- For practice questions, mirror terminology and emphasis found in lecturer-authored materials.
- Do not treat `.url` shortcuts as authoritative content; they are navigation references only.
- Do not assume a video was processed unless a transcript appears in `documents/` or `chunks/`.

## Course organization detected

Detected chapters: **{chapter_text}**.
{f"Potential source gaps: Chapter {', '.join(map(str, gaps))} was not found in downloaded material." if gaps else "No chapter gap can be inferred from detected source numbering."}

Read `START_HERE.md`, `COURSE_MAP.md`, `COVERAGE_REPORT.md`, and `TUTOR_PROTOCOL.md` before teaching. Use `chapters/` as long-form evidence and `meta/corpus.jsonl` for chunk-level retrieval/embeddings.

## Recommended retrieval strategy

For a question about Chapter N:

1. Search chunks where `chapter == N` and priority <= 3.
2. Add `00_course` chunks if the question concerns grading, syllabus, deadlines, or exam scope.
3. Search priority 4 references only if the first two steps are insufficient or the user asks for deeper explanation.
4. Rank chunks with the same relevance by lower numeric `priority` first.

## Folder contract

- `documents/` — normalized full documents in Markdown.
- `chunks/` — chunked Markdown with source metadata and locators.
- `meta/corpus.jsonl` — one JSON object per retrieval chunk; ready for vector DB/RAG import.
- `meta/documents.jsonl` — processing manifest.
- `references_index.md` — reference books intentionally not extracted by default.
- `transcription_queue.md` — media that still needs speech-to-text.
- `links_index.md` — LMS/external shortcuts kept for navigation, not knowledge ingestion.
"""
        write_text(self.output_root / "AI_TUTOR_CONTEXT.md", text)

    def write_reports(self) -> None:
        # References.
        lines = [
            "# References Index",
            "",
            "These files are intentionally kept out of the default AI corpus to avoid flooding retrieval with large books. Use `--include-references` if full extraction is desired.",
            "",
        ]
        if self.reference_paths:
            for x in self.reference_paths:
                lines.append(
                    f"- **{x['title']}** — `{x['source_path']}` ({x['size_mb']} MB)"
                )
        else:
            lines.append("_No reference PDFs were indexed-only._")
        write_text(self.output_root / "references_index.md", "\n".join(lines))

        # Media queue.
        lines = [
            "# Transcription Queue",
            "",
            "Media below was found but not transcribed. Re-run with `--transcribe` to add it to the AI corpus.",
            "",
        ]
        if self.transcription_queue:
            for x in self.transcription_queue:
                lines.append(
                    f"- **{x['title']}** — `{x['source_path']}` ({x['size_mb']} MB)"
                )
        else:
            lines.append("_No pending media._")
        write_text(self.output_root / "transcription_queue.md", "\n".join(lines))

        # URL shortcuts.
        lines = ["# Links Index", ""]
        if self.links:
            for x in self.links:
                url = x["url"] or "(URL not parsed)"
                lines.append(
                    f"- **{x['title']}** — {url} — source `{x['source_path']}`"
                )
        else:
            lines.append("_No .url shortcuts found._")
        write_text(self.output_root / "links_index.md", "\n".join(lines))

        # Errors/skips review.
        duplicate_count = sum(r.status == "duplicate" for r in self.records)
        skipped_html = sum(r.status == "skipped_html_duplicate" for r in self.records)
        skipped_media_stubs = sum(r.status == "skipped_media_stub" for r in self.records)
        lines = [
            "# Processing Report",
            "",
            f"- Ready documents: {sum(r.status == 'ready' for r in self.records)}",
            f"- Chunks: {len(self.chunks)}",
            f"- Duplicate text skipped: {duplicate_count}",
            f"- HTML skipped because content.txt existed: {skipped_html}",
            f"- Tiny media activity wrappers skipped: {skipped_media_stubs}",
            f"- References indexed-only: {len(self.reference_paths)}",
            f"- Media pending transcription: {len(self.transcription_queue)}",
            f"- Errors: {len(self.errors)}",
            "",
        ]
        if self.errors:
            lines.extend(["## Errors", ""])
            for e in self.errors:
                lines.append(f"- `{e['source_path']}` — {e['error']}")
        write_text(self.output_root / "processing_report.md", "\n".join(lines))

    def write_study_pack_outputs(self) -> None:
        self.meta_dir.mkdir(parents=True, exist_ok=True)
        (self.meta_dir / "visual_manifest.json").write_text(
            json.dumps(self.visual_sources, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        records = [asdict(record) for record in self.records]
        write_study_navigation(self.output_root, self.course_name, records)
        validation = validate_ai_study_pack(self.output_root)
        lines = ["# AI Study Pack Quality Report", "", "## Structural errors", ""]
        if validation.errors:
            lines.extend(f"- {error}" for error in validation.errors)
        else:
            lines.append("- None")
        lines.extend(("", "## Warnings", ""))
        if validation.warnings:
            lines.extend(f"- {warning}" for warning in validation.warnings)
        else:
            lines.append("- None")
        write_text(self.output_root / "QUALITY_REPORT.md", "\n".join(lines))
        if validation.errors:
            raise RuntimeError("AI Study Pack validation failed: " + "; ".join(validation.errors))
        self.study_pack_path = create_chatgpt_study_pack(self.output_root, self.course_name)


# -----------------------------------------------------------------------------
# CLI
# -----------------------------------------------------------------------------

def parse_args():
    p = argparse.ArgumentParser(
        description="Convert raw BK-LMS downloads into an AI-ready course knowledge base."
    )
    p.add_argument("--input", required=True, type=Path, help="Raw course folder or .zip")
    p.add_argument("--output", required=True, type=Path, help="Destination knowledge-base folder")

    p.add_argument(
        "--include-references",
        action="store_true",
        help="Extract textbook/reference PDFs into the corpus (default: index only).",
    )

    p.add_argument(
        "--transcribe",
        action="store_true",
        help="Transcribe video/audio using faster-whisper.",
    )
    p.add_argument(
        "--whisper-model",
        default="small",
        help="faster-whisper model name (tiny/base/small/medium/large-v3...). Default: small",
    )
    p.add_argument(
        "--language",
        default=None,
        help="Force transcript language, e.g. vi/en. Default: auto-detect.",
    )
    p.add_argument(
        "--whisper-device",
        default="cpu",
        choices=["cpu", "cuda"],
        help="Whisper device. Default: cpu",
    )
    p.add_argument(
        "--whisper-compute-type",
        default="int8",
        help="faster-whisper compute type. CPU default: int8. CUDA commonly float16.",
    )

    p.add_argument(
        "--chunk-chars",
        type=int,
        default=4800,
        help="Max characters per retrieval chunk. Default: 4800",
    )
    p.add_argument(
        "--chunk-overlap",
        type=int,
        default=500,
        help="Approximate overlap characters. Default: 500",
    )
    p.add_argument(
        "--force",
        action="store_true",
        help="Delete existing output folder before generation.",
    )
    return p.parse_args()


def validate_args(args) -> None:
    if args.chunk_chars < 500:
        raise ValueError("--chunk-chars nên >= 500")
    if args.chunk_overlap < 0:
        raise ValueError("--chunk-overlap phải >= 0")
    if args.chunk_overlap >= args.chunk_chars:
        raise ValueError("--chunk-overlap phải nhỏ hơn --chunk-chars")


def run_preparation(args) -> Path:
    """Run one local knowledge-base build; reusable by the packaged GUI."""
    validate_args(args)

    output_root = args.output.expanduser().resolve()
    if output_root.exists() and args.force:
        shutil.rmtree(output_root)
    output_root.mkdir(parents=True, exist_ok=True)

    # Temporary extraction sits outside output so --force/re-run remains clean.
    with tempfile.TemporaryDirectory(prefix="prepare_ai_course_") as tmp:
        work_dir = Path(tmp)
        source_root, _ = prepare_input(args.input, work_dir)

        print("=" * 78)
        print("PREPARE AI COURSE")
        print("=" * 78)
        print(f"Input : {source_root}")
        print(f"Output: {output_root}")
        print(f"References: {'extract' if args.include_references else 'index only'}")
        print(f"Media transcript: {'yes' if args.transcribe else 'queue only'}")
        print()

        preparer = CoursePreparer(args, source_root, output_root)
        preparer.process()

    print()
    print("=" * 78)
    print("HOÀN TẤT")
    print("=" * 78)
    print(f"Knowledge base: {output_root}")
    print(f"AI context     : {output_root / 'AI_TUTOR_CONTEXT.md'}")
    print(f"Course index   : {output_root / 'course_index.md'}")
    print(f"RAG corpus     : {output_root / 'meta' / 'corpus.jsonl'}")
    print(f"Report         : {output_root / 'processing_report.md'}")
    return output_root


def main():
    run_preparation(parse_args())


if __name__ == "__main__":
    main()
