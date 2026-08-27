from pathlib import Path
import importlib
import sys
import tempfile

# Allow `python app.py` from a source checkout without installing the package.
src_dir = Path(__file__).resolve().parent / "src"
if str(src_dir) not in sys.path:
    sys.path.insert(0, str(src_dir))

from bklms_downloader.gui import main


def _ai_runtime_smoke() -> None:
    """Exercise every packaged AI reader and resource without opening the GUI."""
    from bklms_downloader.ai_prepare import (
        AICoursePreparer,
        REQUIRED_AI_MODULES,
        default_ai_tool_path,
        missing_ai_dependencies,
    )

    missing = missing_ai_dependencies()
    if missing:
        raise RuntimeError("Missing packaged AI modules: " + ", ".join(missing))
    for module_name in REQUIRED_AI_MODULES.values():
        importlib.import_module(module_name)
    if not default_ai_tool_path().is_file():
        raise RuntimeError("Bundled AI preparation tool is missing")

    with tempfile.TemporaryDirectory(prefix="bklms_ai_smoke_") as temp_dir:
        course_root = Path(temp_dir) / "Course"
        course_root.mkdir()
        (course_root / "notes.txt").write_text("AI runtime smoke", encoding="utf-8")

        html_dir = course_root / "Web Page"
        html_dir.mkdir()
        (html_dir / "content.html").write_text(
            "<h1>Packaged HTML smoke</h1><p>Local extraction.</p>",
            encoding="utf-8",
        )

        from pypdf import PdfWriter

        pdf_writer = PdfWriter()
        pdf_writer.add_blank_page(width=72, height=72)
        with (course_root / "tiny.pdf").open("wb") as pdf_handle:
            pdf_writer.write(pdf_handle)

        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Packaged PPTX smoke"
        presentation.save(course_root / "tiny.pptx")

        output = AICoursePreparer().prepare(course_root)
        required_outputs = (
            output / "AI_TUTOR_CONTEXT.md",
            output / "course_index.md",
            output / "processing_report.md",
            output / "meta" / "corpus.jsonl",
        )
        if not all(path.is_file() for path in required_outputs):
            raise RuntimeError("AI runtime self-test did not create required outputs")

if __name__ == "__main__":
    if "--self-test-ai" in sys.argv:
        _ai_runtime_smoke()
    else:
        main()
